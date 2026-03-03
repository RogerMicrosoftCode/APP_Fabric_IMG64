"""
Generador de PowerPoint con Azure OpenAI / Foundry
Genera presentaciones con perfiles de empleados y sus fotos

Requiere:
- pip install python-pptx openai pillow requests azure-identity
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import io
import base64
from PIL import Image
from openai import AzureOpenAI
from azure.identity import DefaultAzureCredential, get_bearer_token_provider
import os
from typing import List, Dict
import json


class EmployeeProfilePPTGenerator:
    """
    Genera presentaciones PowerPoint con perfiles de empleados
    usando Azure OpenAI para generar contenido descriptivo
    """
    
    def __init__(
        self,
        azure_openai_endpoint: str = None,
        azure_openai_key: str = None,
        deployment_name: str = "gpt-4",
        use_managed_identity: bool = True
    ):
        """
        Inicializa el generador
        
        Args:
            azure_openai_endpoint: Endpoint de Azure OpenAI
            azure_openai_key: API Key (opcional si usa Managed Identity)
            deployment_name: Nombre del deployment (gpt-4, gpt-4o, etc.)
            use_managed_identity: Usar DefaultAzureCredential
        """
        self.endpoint = azure_openai_endpoint or os.getenv("AZURE_OPENAI_ENDPOINT")
        self.deployment = deployment_name
        
        if use_managed_identity:
            # Usar Managed Identity / DefaultAzureCredential
            token_provider = get_bearer_token_provider(
                DefaultAzureCredential(),
                "https://cognitiveservices.azure.com/.default"
            )
            self.client = AzureOpenAI(
                azure_endpoint=self.endpoint,
                azure_ad_token_provider=token_provider,
                api_version="2024-02-01"
            )
        else:
            # Usar API Key
            self.client = AzureOpenAI(
                azure_endpoint=self.endpoint,
                api_key=azure_openai_key or os.getenv("AZURE_OPENAI_KEY"),
                api_version="2024-02-01"
            )
    
    
    def generate_profile_description(self, employee_data: Dict) -> str:
        """
        Usa Azure OpenAI para generar descripción profesional del empleado
        
        Args:
            employee_data: Dict con employee_name, department, position
            
        Returns:
            Texto descriptivo generado por IA
        """
        prompt = f"""
Genera una descripción profesional breve (2-3 oraciones) para un perfil de empleado con estos datos:

Nombre: {employee_data.get('employee_name', 'N/A')}
Departamento: {employee_data.get('department', 'N/A')}
Posición: {employee_data.get('position', 'N/A')}

La descripción debe ser positiva, profesional y destacar sus responsabilidades clave.
Escribe en español.
"""
        
        try:
            response = self.client.chat.completions.create(
                model=self.deployment,
                messages=[
                    {"role": "system", "content": "Eres un asistente que genera descripciones profesionales para perfiles de empleados."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=150
            )
            
            return response.choices[0].message.content.strip()
        
        except Exception as e:
            print(f"⚠️ Error generando descripción con OpenAI: {e}")
            return f"Profesional con experiencia en {employee_data.get('department', 'la organización')}."
    
    
    def download_image_from_url(self, url: str) -> Image.Image:
        """
        Descarga imagen desde URL
        
        Args:
            url: URL de la imagen (puede ser HTTP o data URI)
            
        Returns:
            PIL Image
        """
        try:
            if url.startswith('data:image'):
                # Data URI (base64)
                base64_data = url.split(',')[1]
                image_data = base64.b64decode(base64_data)
                return Image.open(io.BytesIO(image_data))
            else:
                # HTTP URL
                response = requests.get(url, timeout=10)
                response.raise_for_status()
                return Image.open(io.BytesIO(response.content))
        
        except Exception as e:
            print(f"❌ Error descargando imagen: {e}")
            # Retornar imagen placeholder
            return Image.new('RGB', (300, 300), color='lightgray')
    
    
    def add_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        """
        Agrega slide de título
        """
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout de título
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
    
    
    def add_employee_slide(
        self,
        prs: Presentation,
        employee_data: Dict,
        generate_ai_description: bool = True
    ):
        """
        Agrega slide con perfil de empleado
        
        Args:
            prs: Presentation object
            employee_data: Dict con:
                - employee_name
                - employee_id
                - department
                - position
                - photo_url (URL HTTP o data URI base64)
            generate_ai_description: Generar descripción con OpenAI
        """
        # Usar layout en blanco
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Título: Nombre del empleado
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = employee_data.get('employee_name', 'Empleado')
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)  # Azul oscuro
        
        # Foto del empleado
        photo_url = employee_data.get('photo_url')
        if photo_url:
            try:
                image = self.download_image_from_url(photo_url)
                
                # Redimensionar y guardar en buffer
                image.thumbnail((400, 400), Image.Resampling.LANCZOS)
                buffer = io.BytesIO()
                image.save(buffer, format='PNG')
                buffer.seek(0)
                
                # Agregar a slide
                left = Inches(1)
                top = Inches(1.8)
                
                slide.shapes.add_picture(buffer, left, top, width=Inches(2.5))
            
            except Exception as e:
                print(f"⚠️ No se pudo agregar foto: {e}")
        
        # Información del empleado
        info_left = Inches(4)
        info_top = Inches(1.8)
        info_width = Inches(5)
        info_height = Inches(1.5)
        
        info_box = slide.shapes.add_textbox(info_left, info_top, info_width, info_height)
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        
        # Posición
        p = info_frame.paragraphs[0]
        p.text = employee_data.get('position', 'N/A')
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(68, 68, 68)
        
        # Departamento
        p = info_frame.add_paragraph()
        p.text = f"📍 {employee_data.get('department', 'N/A')}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)
        
        # Employee ID
        p = info_frame.add_paragraph()
        p.text = f"🆔 ID: {employee_data.get('employee_id', 'N/A')}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.space_before = Pt(6)
        
        # Descripción generada por IA
        if generate_ai_description:
            description = self.generate_profile_description(employee_data)
        else:
            description = employee_data.get('description', '')
        
        if description:
            desc_top = Inches(3.8)
            desc_height = Inches(1.5)
            
            desc_box = slide.shapes.add_textbox(info_left, desc_top, info_width, desc_height)
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            
            p = desc_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(14)
            p.line_spacing = 1.2
    
    
    def generate_presentation(
        self,
        employees: List[Dict],
        output_path: str = "employee_profiles.pptx",
        title: str = "Perfiles de Empleados",
        subtitle: str = "",
        use_ai_descriptions: bool = True
    ) -> str:
        """
        Genera presentación completa
        
        Args:
            employees: Lista de dicts con datos de empleados
            output_path: Ruta para guardar el .pptx
            title: Título de la presentación
            subtitle: Subtítulo
            use_ai_descriptions: Generar descripciones con OpenAI
            
        Returns:
            Path del archivo generado
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide de título
        self.add_title_slide(prs, title, subtitle)
        
        # Slides de empleados
        for i, emp in enumerate(employees):
            print(f"📄 Generando slide {i+1}/{len(employees)}: {emp.get('employee_name', 'N/A')}")
            self.add_employee_slide(prs, emp, generate_ai_description=use_ai_descriptions)
        
        # Guardar
        prs.save(output_path)
        print(f"✅ Presentación guardada: {output_path}")
        
        return output_path


# ============================================================================
# EJEMPLO DE USO
# ============================================================================

def main():
    """
    Ejemplo de uso del generador
    """
    
    # Configurar generador
    generator = EmployeeProfilePPTGenerator(
        azure_openai_endpoint="https://your-openai.openai.azure.com/",
        deployment_name="gpt-4o",
        use_managed_identity=True
    )
    
    # Datos de empleados (normalmente vendrían de Fabric Data Agent)
    employees = [
        {
            "employee_id": "12345",
            "employee_name": "Pamela Yissell",
            "department": "Human Resources",
            "position": "Senior HR Manager",
            "photo_url": "https://your-api.azurecontainerapps.io/image/onelake/12345"
        },
        {
            "employee_id": "67890",
            "employee_name": "Juan Pérez",
            "department": "Engineering",
            "position": "Software Engineer",
            "photo_url": "https://your-api.azurecontainerapps.io/image/onelake/67890"
        }
    ]
    
    # Generar presentación
    output_file = generator.generate_presentation(
        employees=employees,
        output_path="perfiles_empleados.pptx",
        title="Directorio de Empleados",
        subtitle="Generado automáticamente con Azure OpenAI",
        use_ai_descriptions=True
    )
    
    print(f"\n✅ Presentación generada exitosamente: {output_file}")


def create_from_fabric_agent_response(agent_response_json: str, output_path: str):
    """
    Crea PowerPoint desde respuesta JSON del Fabric Data Agent
    
    Args:
        agent_response_json: JSON string con datos del agente
        output_path: Ruta para guardar el .pptx
    """
    
    # Parse JSON
    data = json.loads(agent_response_json)
    
    # Extraer empleados (ajustar según estructura del agente)
    employees = []
    for row in data.get('rows', []):
        employees.append({
            'employee_id': row.get('employee_id'),
            'employee_name': row.get('employee_name'),
            'department': row.get('department'),
            'position': row.get('position'),
            'photo_url': row.get('photo_url')
        })
    
    # Generar PPT
    generator = EmployeeProfilePPTGenerator(use_managed_identity=True)
    generator.generate_presentation(
        employees=employees,
        output_path=output_path,
        title="Perfiles Consultados",
        use_ai_descriptions=True
    )


if __name__ == "__main__":
    main()


"""
INTEGRACIÓN CON COPILOT STUDIO
================================

1. Crear Skill/Action en Copilot Studio:
   
   Nombre: GenerateEmployeePPT
   Type: HTTP Request
   URL: https://your-function-app.azurewebsites.net/api/generate-ppt
   Method: POST
   Body: {
       "employees": [
           {"employee_id": "12345", "employee_name": "...", ...}
       ],
       "title": "Mi Presentación"
   }

2. Desde el Topic de Copilot:
   
   - Query Fabric Data Agent para obtener empleados
   - Parsear respuesta JSON
   - Llamar a GenerateEmployeePPT Action
   - Retornar enlace de descarga al usuario

3. Deploy como Azure Function:
   
   func init EmployeePPTFunction --python
   # Agregar generate_ppt_with_openai.py como módulo
   func azure functionapp publish YOUR_FUNCTION_APP

4. Variables de entorno necesarias:
   
   AZURE_OPENAI_ENDPOINT=https://your-openai.openai.azure.com/
   AZURE_OPENAI_DEPLOYMENT=gpt-4o
   FABRIC_API_ENDPOINT=https://api.fabric.microsoft.com
"""
